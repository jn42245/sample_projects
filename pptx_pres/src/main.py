# Created on 5 December 2022
# Created by Thierry Tran
# The purpose of the script is to generate automatic PowerPoint presentations from Excel dataset

import os
import pandas as pd
from pptx import Presentation, util


def title_slide(presentation: Presentation, title_str: str, placeholder_str: str, layout: int = 0):
    s = presentation.slides.add_slide(presentation.slide_layouts[layout])
    s.shapes.title.text = title_str
    s.placeholders[1].text = placeholder_str
    return s


def image_text_slide(presentation: Presentation, title_str: str, fig_path: str, fig_x: float,
                     fig_y: float, fig_width: float, fig_height: float, text_x: float,
                     text_y: float, text_width: float, text_height: float, text_display: list,
                     paragraph_spacing: int, layout: int = 3):
    s = presentation.slides.add_slide(presentation.slide_layouts[layout])
    s.shapes.title.text = title_str

    s.shapes.add_picture(fig_path, util.Inches(fig_x), util.Inches(fig_y),
                         width=util.Inches(fig_width), height=util.Inches(fig_height))

    os.remove(fig_path)

    text_box = s.shapes.add_textbox(util.Inches(text_x), util.Inches(text_y),
                                    width=util.Inches(text_width), height=util.Inches(text_height))
    tf = text_box.text_frame
    tf.word_wrap = True

    for x in range(0, len(text_display)):
        p = tf.add_paragraph()
        if 0 < x:
            p.space_before = paragraph_spacing
        p.text = text_display[x]
    return s


def table_slide(presentation: Presentation, title_str: str,
                table_x: float, table_y: float, c_x: float, c_y: float,
                table_style: str, rounding: int, raw_excel: bool =True, df: pd.DataFrame = pd.DataFrame(),
                excel_path: str = '', sheet: str = '', layout: int = 3):

    if raw_excel:
        slide_table = pd.read_excel(excel_path, sheet_name=sheet, header=None)
    else:
        slide_table = df

    s = presentation.slides.add_slide(presentation.slide_layouts[layout])
    s.shapes.title.text = title_str

    shape = s.shapes.add_table(len(slide_table.index), len(slide_table.columns),
                               util.Inches(table_x), util.Inches(table_y), util.Inches(c_x), util.Inches(c_y))

    table = shape.table
    tbl = shape._element.graphic.graphicData.tbl
    tbl[0][-1].text = '{' + table_style + '}'

    for row in range(0, len(slide_table.index)):
        for col in range(0, len(slide_table.columns)):
            if isinstance(slide_table.iat[row, col], float) or isinstance(slide_table.iat[row, col], int):
                table.cell(row, col). text = str(round(slide_table.iat[row, col], rounding))
            else:
                table.cell(row, col).text = str(slide_table.iat[row, col])
    return s


def transition_slide(presentation: Presentation, placeholder_str: str, layout: int = 10):
    s = presentation.slides.add_slide(presentation.slide_layouts[layout])
    s.placeholders[0].text = placeholder_str
    return s


def final_slde(presentation: Presentation, layout: int = 22):
    s = presentation.slides.add_slide(presentation.slide_layouts[layout])
    return s


def remove_slide(presentation: Presentation, lower_bound: int, upper_bound: int):
    xml_slides = presentation.slides._sldIdLst
    slide_count = list(xml_slides)

    for x in range(lower_bound, upper_bound):
        xml_slides.remove(slide_count[x])
    return presentation


if __name__ == '__main__':
    import os
    import calendar
    import logging
    import pandas as pd
    from datetime import datetime
    from pathlib import Path

    # Internal module
    from src.dataload.json_objects import load_json
    from src.utils.plotly_graphs import pie_chart, bar_chart, stacked_bar_chart, scatter_plot
    from src.utils.maps import highlight_countries_world
    from src.utils.text_generator import compare_data_text, countries_highlighted_text

    # Variables needed for the presentation
    wk_dir = os.path.dirname(os.path.abspath('__file__'))
    sqt_colours = load_json(os.path.join(wk_dir, 'data', 'sqt_colours.json'), "r")
    table_styles = load_json(os.path.join(wk_dir, 'data', 'table_styles.json'), "r")
    date_str = datetime.today()

    logging.basicConfig(level=logging.INFO, format="%(asctime)s - %(levelname)s - %(message)s")
    
    logging.info('Reading datafile')
    
    table_df = pd.read_excel(os.path.join(wk_dir, 'data', 'sample_table_additional_data.xlsx'), sheet_name='sample', index_col=0)
    for row in range(0, len(table_df.index)):
        for col in range(0, len(table_df.columns)):
            if isinstance(table_df.iat[row, col], float):
                table_df.iat[row, col] = int(table_df.iat[row, col])

    logging.info('Creating presentation')
                
    pres = Presentation(os.path.join(wk_dir, 'data', 'sqt_template.pptx'))

    # Title slide
    title_slide(pres, 'Example presentation', '{d} {m} {y}'.format(
        d=date_str.day, m=calendar.month_name[date_str.month], y=date_str.year))

    # Transition slide
    transition_slide(pres, 'EXAMPLE - Slide with charts and automatic text')

    # Slide with image and text
    # Create example graph
    labels = [*table_df.columns[1:]]
    dict_values_multiple = {}
    dict_colours_multiple = {}

    for ix in table_df.index:
        dict_values_multiple[ix] = table_df.loc[table_df.index == ix, labels].values.flatten().tolist()

    for ix in range(0, len(dict_values_multiple)):
        dict_colours_multiple[list(dict_values_multiple.keys())[ix]] = [sqt_colours[x] for x in sqt_colours.keys()][ix]

    fig = bar_chart(labels, dict_values_multiple, dict_colours_multiple, 1500, 1000)
    fig.write_image(os.path.join(wk_dir, 'output', 'temp.jpeg'))

    # Example text
    paragraphs = []
    for key in dict_values_multiple:
        paragraphs.append(compare_data_text(dict_values_multiple[key][0],
                                            dict_values_multiple[key][-1], labels[0], labels[-1],
                                            key, 'sales'))

    # Creating the slide
    image_text_slide(pres, 'EXAMPLE - Slide with charts and automatic text',
                     os.path.join(wk_dir, 'output', 'temp.jpeg'), 0.5,
                     1.5, 4.5, 3, 5.5, 1.5, 4, 3, paragraphs, 76200)

    # Transition slide
    transition_slide(pres, 'EXAMPLE - Slide with table')

    # Slide with Table
    table_slide(pres, 'EXAMPLE - Slide with table', table_x=1, table_y=1, c_x=8, c_y=1,
                table_style=table_styles["LightStyle2"], rounding=0, raw_excel=False,
                df=table_df.reset_index().T.reset_index().T)

    # Transition slide
    transition_slide(pres, 'EXAMPLE - Slides with maps')

    # Slide with world map
    highlight_countries_world('gainsboro', 'black', 1, dict_colours_multiple,
                              Path(os.path.join(wk_dir, 'output', 'temp_map.jpeg')), cities=True)

    image_text_slide(pres, 'EXAMPLE - Slide with world map',
                     os.path.join(wk_dir, 'output', 'temp_map.jpeg'),
                     0.5, 0.7, 4.5, 4.5, 5.5, 1.5, 4, 3,
                     [countries_highlighted_text(list(dict_colours_multiple.keys()))], 76200)

    # Region slides
    for cont in table_df['continent'].unique().tolist():
        temp_df = table_df[table_df['continent'] == cont].copy()
        temp_dict_values_multiple = {}
        temp_dict_colours_multiple = {}

        for ix in temp_df.index:
            temp_dict_values_multiple[ix] = temp_df.loc[temp_df.index == ix, labels].values.flatten().tolist()

        for ix in range(0, len(temp_dict_values_multiple)):
            temp_dict_colours_multiple[list(
                temp_dict_values_multiple.keys())[ix]] = [sqt_colours[x] for x in sqt_colours.keys()][ix]

        highlight_countries_world('gainsboro', 'black', 1, temp_dict_colours_multiple,
                                  Path(os.path.join(wk_dir, 'output', 'temp_map.jpeg')), continent=True,
                                  continent_str=cont)

        image_text_slide(pres, 'EXAMPLE - Slide with {} map'.format(cont),
                         os.path.join(wk_dir, 'output', 'temp_map.jpeg'),
                         0.5, 0.7, 4.5, 4.5, 5.5, 1.5, 4, 3,
                         [countries_highlighted_text(list(temp_dict_colours_multiple.keys()))], 76200)

    # Final slide
    final_slde(pres)

    # Remove additional slides in template
    pres = remove_slide(pres, 0, 96)

    pres.save(os.path.join(wk_dir, 'output', 'example_presentation.pptx'))
    
    logging.info('Presentation created and saved')
